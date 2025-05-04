import os
from sentence_transformers import SentenceTransformer
import PyPDF2
import docx2txt
from pptx import Presentation
from sklearn.metrics.pairwise import cosine_similarity
from langchain_text_splitters import RecursiveCharacterTextSplitter
import numpy as np
from config import EMBEDDING_MODEL, CHUNK_SIZE, CHUNK_OVERLAP, DEFAULT_ENCODING, SUPPORTED_FILE_TYPES

class DocumentParser:
    def __init__(self):
        self.embedder = SentenceTransformer(EMBEDDING_MODEL)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP
        )

    def extract_text(self, file_path):
        """Extract text from various file types"""
        file_ext = os.path.splitext(file_path.lower())[1]
        if file_ext not in SUPPORTED_FILE_TYPES:
            raise ValueError(f"Unsupported file type: {file_ext}. Supported types: {SUPPORTED_FILE_TYPES}")

        if file_ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = " ".join([page.extract_text() for page in reader.pages])
        elif file_ext == '.docx':
            text = docx2txt.process(file_path)
        elif file_ext == '.pptx':
            prs = Presentation(file_path)
            text = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_ext == '.txt':
            with open(file_path, 'r', encoding=DEFAULT_ENCODING) as f:
                text = f.read()
        
        return text.strip()

    def process_document(self, file_path):
        """
        Process a document file and return chunks and embeddings.
        """
        try:
            # Extract text
            text = self.extract_text(file_path)
            
            # Split into chunks
            chunks = self.text_splitter.split_text(text)
            
            # Generate embeddings
            embeddings = self.embedder.encode(chunks)
            
            return {
                'chunks': chunks,
                'embeddings': embeddings
            }
        except Exception as e:
            print(f"Error processing document: {str(e)}")
            return None

    def process_multiple_documents(self, file_paths):
        """
        Process multiple documents and combine their chunks and embeddings.
        """
        all_chunks = []
        all_embeddings = []
        
        for file_path in file_paths:
            result = self.process_document(file_path)
            if result:
                all_chunks.extend(result['chunks'])
                all_embeddings.append(result['embeddings'])
        
        if not all_chunks:
            return None
            
        # Combine embeddings
        combined_embeddings = np.vstack(all_embeddings)
        
        return {
            'chunks': all_chunks,
            'embeddings': combined_embeddings
        }
